from sitesyncro import Model

from sitesyncro.utils.fnc_visualize import (pygraphviz_layout)
from sitesyncro.utils.fnc_phase import (reduce_earlier_than)

import os
import shutil
import copy
import io
from collections import defaultdict

import numpy as np
import networkx as nx
import matplotlib.pyplot as pyplot
import matplotlib.animation as animation
from matplotlib.animation import PillowWriter
from matplotlib.lines import Line2D
from natsort import natsorted
from PIL import Image, ImageSequence
from pptx import Presentation
from pptx.util import Inches


def convert(old_filename, new_filename, duration):
	images = []
	with Image.open(old_filename) as im:
		for i in range(im.n_frames):
			im.seek(i)
			buf = io.BytesIO()
			im.save(buf, format='png')
			buf.seek(0)
			images.append(Image.open(buf))
	images[0].save(new_filename, save_all=True, append_images=images[1:], optimize=False, duration=duration)


def save_as_ppt(fname, img_paths, slide_size = (13.33, 7.5)):
	
	prs = Presentation()
	slide_layout = prs.slide_layouts[6]  # Using a blank layout
	prs.slide_width = Inches(slide_size[0])
	prs.slide_height = Inches(slide_size[1])
	slide_width = prs.slide_width
	slide_height = prs.slide_height
	for img_path in img_paths:
		img = Image.open(img_path)
		img_width, img_height = img.size
		width_ratio = slide_width / img_width
		height_ratio = slide_height / img_height
		fit_ratio = min(width_ratio, height_ratio)
		new_width = int(img_width * fit_ratio)
		new_height = int(img_height * fit_ratio)
		left = (slide_width - new_width) / 2
		top = (slide_height - new_height) / 2
		slide = prs.slides.add_slide(slide_layout)
		left = 0
		top = 0
		pic = slide.shapes.add_picture(img_path, left, top, new_width, new_height)
	prs.save(fname)


def get_graph(earlier_than, samples, gap = 0.05):
	
	def _update_pos_by_groups(pos, groups, gap = 0.05):
		
		x_range = max(x for x, _ in pos.values()) - min(x for x, _ in pos.values())
		gap *= x_range
		
		x_last = 0
		for group in groups:
			x_min = min(pos[n][0] for n in group)
			for n in group:
				pos[n] = ((pos[n][0] - x_min) + x_last, pos[n][1])
			x_last = max(pos[n][0] for n in group) + gap
		
		x_last = 0
		for n in sorted(pos.keys(), key = lambda n: pos[n][0]):
			if pos[n][0] < x_last + gap:
				pos[n] = (x_last + gap, pos[n][1])
				x_last = pos[n][0]
		
		return pos
		
	earlier_than = reduce_earlier_than(earlier_than)
	G = nx.convert_matrix.from_numpy_array(earlier_than, create_using=nx.DiGraph)
	pos = pygraphviz_layout(G)  # {i: [x,y], ...}
	
	done = set()
	groups = []
	for c in nx.connected_components(G.to_undirected()):
		G_sub = G.subgraph(c)
		groups.append(list(G_sub.nodes()))
		done.update(groups[-1])
	for n in pos:
		if n not in done:
			groups.append([n])
	
	groups = natsorted(groups, key = lambda group: samples[group[0]])
	pos = _update_pos_by_groups(pos, groups)
	
	groups = dict([(n, i) for i, group in enumerate(groups) for n in group])
	
	return G, pos, groups, gap


def get_G(earlier_than):
	
	earlier_than = reduce_earlier_than(earlier_than)
	return nx.convert_matrix.from_numpy_array(earlier_than, create_using=nx.DiGraph)


def update_pos_by_dists(pos, samples, dists):
	pos = copy.deepcopy(pos)
	t_range = max(dists[s][0] for s in samples) - min(dists[s][0] for s in samples)
	height = max(y for x, y in pos.values()) - min(y for x, y in pos.values())
	for i, s in enumerate(samples):
		pos[i] = (pos[i][0] * height / t_range, dists[s][0])
	return pos


def interpolate_dists(dists1, dists2, t):
	dists = {}
	for s in dists1:
		dists[s] = ((1 - t) * np.array(dists1[s]) + t * np.array(dists2[s])).tolist()
	return dists


def update_graph(i, ax, title, G, pos, node_color, edge_color, likelihoods, posteriors, samples, outliers, phases, clusters, t_min, t_max, t_step, prev_posteriors, total_frames):
	
	if posteriors:
		if prev_posteriors is not None:
			pos_final = update_pos_by_dists(pos, samples, posteriors)
			posteriors = interpolate_dists(prev_posteriors, posteriors, i / total_frames)
			pos = update_pos_by_dists(pos, samples, posteriors)
		else:
			pos = update_pos_by_dists(pos, samples, posteriors)
			pos_final = pos
	else:
		pos_final = pos
	
	handles_clusters = []
	handles_dates = []
	ax.clear()
	ax.set_title(title)
	x_values = [pos[i][0] for i in range(len(samples))]
	x_min = min(x_values)
	x_max = max(x_values)
	step = (x_max - x_min) / len(pos)
	if posteriors:
		graph_width = x_max - x_min
		whisker_length = graph_width * 0.005
		for i, s in enumerate(samples):
			x = pos[i][0]
			
			lbl_posterior, lbl_likelihood = None, None
			if i == len(samples) - 1:
				lbl_posterior = "Modelled"
				lbl_likelihood = "Unmodelled"
			
			color = "k"
			if i in outliers:
				color = "r"
			m, r1, r2 = likelihoods[s]
			handle = ax.errorbar(x, m, yerr=[[m-r2],[r1-m]], fmt='.', color=color, markersize=5, linewidth=2, zorder=1, alpha=0.15, label=lbl_likelihood)
			if color == 'k':
				handle_likelihood = handle
			ax.hlines([r1, r2], x - whisker_length, x + whisker_length, color=color, linewidth=2, zorder=1, alpha=0.15)
			
			m, r1, r2 = posteriors[s]
			handle_posterior = ax.errorbar(x, m, yerr=[[m-r2],[r1-m]], fmt='.', color='k', markersize=5, linewidth=.5, zorder=2, label=lbl_posterior)
			ax.hlines([r1, r2], x - whisker_length, x + whisker_length, color='k', linewidth=.5, zorder=2)
		handles_dates.append(handle_likelihood)
		handles_dates.append(handle_posterior)
		
		t0 = np.floor((1950 - t_max) / t_step) * t_step
		t1 = np.ceil((1950 - t_min) / t_step) * t_step
		yticks = np.arange(t0, t1 + t_step, t_step)
		yticks = list(zip(1950 - yticks, yticks.astype(int)))
		ax.set_ylabel("Year CE", labelpad=25)
		ax.set_ylim(t_min - t_step, t_max + t_step)

	else:
		yticks = sorted(list(set([y for _, y in pos.values()])))
		yticks = list(zip(yticks, np.arange(len(yticks), dtype=int)[::-1] + 1))
		ax.set_ylabel("Stratigraphic Phase", rotation = 90, labelpad=10)

	nodes = nx.draw_networkx_nodes(G, pos, node_color=node_color, node_size=25, ax=ax)
	nx.draw_networkx_edges(G, pos, edge_color=edge_color, arrows=True, width=.5, ax=ax)
	nodes.set_zorder(3)

	y_min = ax.get_ylim()[0]
	y_max = ax.get_ylim()[1]
	gap = 0.005 * (ax.get_ylim()[1] - y_min)
	if not phases:

		group_lines = []
		last_group = None
		for n in sorted(pos.keys(), key=lambda n: pos[n][0]):
			if groups[n] != last_group:
				group_lines.append(pos[n][0] - step / 2)
				last_group = groups[n]
		group_lines.append(pos[n][0])
		group_lines = group_lines[1:]

		last_x = x_min
		for i, x in enumerate(group_lines):
			prefix = "Group" if ((x - last_x) / step) > 1.5 else "Gr."
			ax.text((x + last_x) / 2, y_min + gap, "%s %d" % (prefix, i+1), verticalalignment='top', horizontalalignment='center', fontsize=10)
			if i < len(group_lines) - 1:
				ax.axvline(x, color='grey', linewidth=1, linestyle='--')
			last_x = x
	
	for i, s in enumerate(samples):
		ax.text(pos[i][0], y_max + gap, s.split("_")[0], rotation=90, verticalalignment='top', horizontalalignment='center', fontsize=10)
	
	for y, label in yticks:
		ax.text(step / 2, y, str(label), verticalalignment='top', horizontalalignment='right', fontsize=8)

	if phases:
		phase_lines = []
		last_phase = None
		last_y = 0
		labels = []
		for n in sorted(pos_final.keys(), key=lambda n: pos_final[n][1]):
			if phases[n] != last_phase:
				phase_lines.append((pos_final[n][1] + last_y)/2)
				last_phase = phases[n]
				labels.append(last_phase)
			last_y = pos_final[n][1]
		phase_lines.append((pos_final[n][1] + last_y)/2)
		phase_lines = phase_lines[1:]
		phase_lines[-1] += 2*step

		phase_lines = [t for t in phase_lines if t > t_min]

		for i, y in enumerate(phase_lines):
			ax.axhline(y, color='grey', linewidth=1, linestyle='--')

		last_y = y_min
		for i, y in enumerate(phase_lines):
			ax.text(x_max + 0.6*step, (last_y + y)/2, "Phase %d" % (labels[i]), verticalalignment='center', horizontalalignment='left', fontsize=10)
			last_y = y

	if clusters:
		collect = defaultdict(list)
		means = defaultdict(list)
		for n in clusters:
			collect[clusters[n]].append(n)
			means[clusters[n]].append(pos_final[n][1])
		means = dict([(label, np.median(means[label])) for label in means])
		labels = sorted(means.keys(), key = lambda label: means[label], reverse=True)
		clusters = dict([(n+1, collect[label]) for n, label in enumerate(labels)])
		means = dict([(n+1, means[label]) for n, label in enumerate(labels)])
		for label in clusters:
			xy = np.array([pos[n] for n in clusters[label]])
			handle, = ax.plot(xy[:,0], xy[:,1], "o", label=str(label), zorder = 3)
			handles_clusters.append(handle)
	
	
	legend1 = None
	if handles_clusters:
		legend1 = ax.legend(title='Cluster', handles=handles_clusters, loc='upper right')
	
	handles_edges = []
	if len(G.edges()) > 0:
		if isinstance(edge_color, list):
			colors = sorted(list(set(edge_color)), key = lambda color: ['k','r','lightgrey'].index(color))
		else:
			colors = [edge_color]
		if colors:
			handles_edges.append(Line2D([], [], linestyle='None', label='Chronological Relations', color='k'))
		for color in colors:
			label = 'Earlier-than by stratigraphy' if color == 'k' else 'Earlier-than by dating'
			handle = Line2D([0], [0], color=color, lw=1, linestyle='-', marker='>', markersize=5, label=label)
			handles_edges.append(handle)
	
	if handles_dates:
		handles_edges.append(Line2D([], [], linestyle='None', label='Dating ranges (95.45%)', color='k'))
		handles_edges += handles_dates
	
	legend2 = ax.legend(handles=handles_edges, loc='lower left')
	for text in legend2.get_texts():
		if text.get_text() in ['Chronological Relations', 'Dating ranges (95.45%)']:
			text.set_position((-40, 0))
	
	if legend1:
		ax.add_artist(legend1)
	ax.add_artist(legend2)
	
	ax.set_xlim(x_min - step/2, x_max + step/2)
	ax.invert_yaxis()
	pyplot.tight_layout()
	pyplot.subplots_adjust(left=0.03, right=0.96)


def plot_graph(title, G, pos, node_color, edge_color, likelihoods, posteriors, samples, outliers, phases, clusters, directory, name, t_min=None, t_max=None, t_step=100, prev_posteriors=None, total_frames=30):
	
	print("Plotting", name, title)
	if prev_posteriors is not None:
		fig, ax = pyplot.subplots(figsize=(20, 10))
		ani = animation.FuncAnimation(fig, update_graph, frames=total_frames, fargs=(ax, title, G, pos, node_color, edge_color, likelihoods, posteriors, samples, outliers, phases, clusters, t_min, t_max, t_step, prev_posteriors, total_frames), repeat=False)
		writer = PillowWriter(fps=25)
		fname = os.path.join(directory, "%s.gif" % name)
		ani.save(fname, writer=writer)
		convert(fname, fname, 120)
	else:
		fig, ax = pyplot.subplots(figsize=(20, 10))
		update_graph(0, ax, title, G, pos, node_color, edge_color, likelihoods, posteriors, samples, outliers, phases, clusters, t_min, t_max, t_step, prev_posteriors, total_frames)
		fname = os.path.join(directory, "%s.png" % name)
		pyplot.savefig(fname)
		pyplot.close()
	return fname


#T_MIN = -1400
T_MIN = None
#T_MAX = -2200
T_MAX = None

DIRECTORY = "vis_stages"

if __name__ == '__main__':
	
	print("Loading data")
	model0 = Model(directory="stage_0")
	model1 = Model(directory="stage_1")
	model2 = Model(directory="stage_2")
	model3 = Model(directory="stage_3")
	
	samples = list(model0.samples.keys())
	
	eaps = dict([(i, model0.samples[s].excavation_area_phase) for i, s in enumerate(samples)])
	phases2 = dict([(i, model2.samples[s].phase) for i, s in enumerate(samples)])
	phases3 = dict([(i, model3.samples[s].phase) for i, s in enumerate(samples)])
	
	outliers = [samples.index(s) for s in model1.outliers]
	data = model2.clusters[model2.cluster_opt_n]
	clusters = {}
	for label in data:
		for s in data[label]:
			clusters[samples.index(s)] = label
	
	print("Getting ranges")
	dists0 = {}
	for name in samples:
		m = model0.samples[name].likelihood_mean
		r1, r2 = model0.samples[name].likelihood_range
		dists0[name] = [m, r1, r2]
	dists1 = {}
	for name in samples:
		m = model1.samples[name].posterior_mean
		r1, r2 = model1.samples[name].posterior_range
		dists1[name] = [m, r1, r2]
	dists2 = {}
	for name in samples:
		m = model2.samples[name].posterior_mean
		r1, r2 = model2.samples[name].posterior_range
		dists2[name] = [m, r1, r2]
	dists3 = {}
	for name in samples:
		m = model3.samples[name].posterior_mean
		r1, r2 = model3.samples[name].posterior_range
		dists3[name] = [m, r1, r2]
	
	print("Populating graph")
	earlier_than0, samples0 = model0.mphasing.create_earlier_than_matrix()
	assert(samples == samples0)
	
	earlier_than1, samples1 = model1.mphasing.create_earlier_than_matrix()
	assert(samples == samples1)
	
	earlier_than2 = model1.mphasing.update_earlier_than_by_dating(earlier_than1, samples)
	
	earlier_than3, samples3 = model3.mphasing.create_earlier_than_matrix()
	assert(samples == samples3)
	
	G0, pos, groups, gap = get_graph(earlier_than0, samples)
	G1 = get_G(earlier_than1)
	G3 = get_G(np.zeros(earlier_than0.shape, dtype = bool))
	
	node_colors0 = 'k'
	node_colors_o = [('r' if i in outliers else 'k') for i in range(len(samples))]
	
	by_dating = []
	et_by_dating = np.zeros(earlier_than0.shape, dtype = bool)
	G2 = get_G(et_by_dating)
	for n in sorted(pos.keys(), key = lambda n: pos[n][0]):
		et_by_dating_s = np.zeros(earlier_than0.shape, dtype = bool)
		for i, j in zip(*np.where(earlier_than2)):
			if i != n:
				continue
			if (not earlier_than1[i, j]) and (groups[i] != groups[j]):
				et_by_dating_s[i,j] = True
		if not et_by_dating_s.any():
			continue
		G_by_dating = get_G(et_by_dating_s)
		for i, j in G_by_dating.edges():
			G2.add_edge(i, j)
		edge_colors_s = []
		for i, j in G2.edges():
			if et_by_dating_s[i, j]:
				color = 'r'
			elif et_by_dating[i, j]:
				color = 'lightgrey'
			edge_colors_s.append(color)
		by_dating.append([samples[n].split("_")[0], G2.copy(), edge_colors_s])
		et_by_dating = et_by_dating | et_by_dating_s
	
	t_min, t_max = np.inf, -np.inf
	for dists in [dists0, dists1, dists2, dists3]:
		for i, s in enumerate(samples):
			m, r1, r2 = dists[s]
			t_min = min(t_min, r2)
			t_max = max(t_max, r1)
	if T_MIN is not None:
		t_min = max(t_min, 1950 - T_MIN)
	if T_MAX is not None:
		t_max = min(t_max, 1950 - T_MAX)
	
	DIRECTORY
	if os.path.isdir(DIRECTORY):
		shutil.rmtree(DIRECTORY)
	os.makedirs(DIRECTORY)
	
	img_paths = [
		plot_graph(	"Stage 1 - Stratigraphic Phasing",
			G0, pos, 'k', 			'k', {}, 	  {},	  samples, [], 		 None, None, DIRECTORY, '01-stage1a'),
		plot_graph(	"Stage 1 - Outlier Detection",
			G1, pos, 'k', 			'k',  dists0, dists0, samples, [], 		 None, None, DIRECTORY, '02-stage1b', t_min, t_max),
		plot_graph(	"Stage 1 - Outlier Detection",
			G1, pos, node_colors_o,	'k', dists0,  dists0, samples, outliers, None, None, DIRECTORY, '03-stage1c', t_min, t_max),
		plot_graph(	"Stage 1 - Chronological Modeling 1",
			G1, pos, 'k', 			'k', dists0,  dists1, samples, outliers, None, None, DIRECTORY, '04-stage1d', t_min, t_max, prev_posteriors=dists0),
	]
	for i, (name, G2_g, edge_colors_g) in enumerate(by_dating):
		img_paths.append(
			plot_graph(	"Stage 2 - Inter-Group Chronological Relations - %s" % (name),
				G2_g, pos, 'k', edge_colors_g, dists0, dists1, samples, outliers, None, None, DIRECTORY, '05-stage2a-%03d' % (i+1), t_min, t_max
			)
		)
	img_paths += [	
		plot_graph(	"Stage 2 - Chronological Modeling 2",
		 	G2, pos, 'k', 'lightgrey',	 dists0, dists2, samples, outliers, phases2, None,	   DIRECTORY, '07-stage2c', t_min, t_max, prev_posteriors=dists1),
		plot_graph(	"Stage 3 - Chronological Clustering",
		 	G3, pos, 'k', 'k', 			 dists0, dists2, samples, outliers, phases2, clusters, DIRECTORY, '08-stage3a', t_min, t_max),
		plot_graph(	"Stage 3 - Chronological Modeling 3",
		 	G3, pos, 'k', 'k', 			 dists0, dists3, samples, outliers, phases3, clusters, DIRECTORY, '09-stage3b', t_min, t_max, prev_posteriors=dists2),
	]
	save_as_ppt(os.path.join(DIRECTORY, "model.pptx"), img_paths)
	